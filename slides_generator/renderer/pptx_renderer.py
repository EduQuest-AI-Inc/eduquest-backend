"""
PPTX Renderer

Converts a list of `CompletedSlide` objects (from the orchestrator) into a
.pptx file using python-pptx. All layout logic lives here; no AI calls happen
at this stage.

Supported layouts:
  title         — Opening slide: lesson name, period info
  concept_intro — Concept name + bullets + prerequisites badge
  two_col       — Bullets left, image right
  visual_focus  — Large image top, caption below
  skill_card    — Skill name, Bloom badge, difficulty badge, mastery bullets
  summary       — Recap of all skills
"""

from __future__ import annotations

import io
import os
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from slides_generator.models.slide_plan import CompletedSlide

# ── Brand colours ────────────────────────────────────────────────────────────
_BLUE = RGBColor(0x1B, 0x4F, 0x9B)
_ORANGE = RGBColor(0xF0, 0x7B, 0x3F)
_DARK = RGBColor(0x2C, 0x2C, 0x2C)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
_PLACEHOLDER_BG = RGBColor(0xDD, 0xDD, 0xDD)
_MUTED = RGBColor(0x88, 0x88, 0x88)

# Bloom level badge colours
_BLOOM_COLOURS: dict[str, RGBColor] = {
    "remember": RGBColor(0x78, 0x90, 0x9C),
    "understand": RGBColor(0x43, 0xA0, 0x47),
    "apply": RGBColor(0x1E, 0x88, 0xE5),
    "analyze": RGBColor(0x8E, 0x24, 0xAA),
    "evaluate": RGBColor(0xF4, 0x51, 0x1E),
    "create": RGBColor(0xC6, 0x28, 0x28),
}

# Slide dimensions (16:9 widescreen)
_W = Inches(13.33)
_H = Inches(7.5)
_MARGIN = Inches(0.5)
_CONTENT_TOP = Inches(1.6)


def render(slides: List[CompletedSlide], meta: Optional[dict] = None) -> bytes:
    """Render a list of `CompletedSlide` into a .pptx file.

    Args:
        slides: Ordered list of CompletedSlide objects from the orchestrator.
        meta: Optional dict with keys lesson_name, period_name, grade_level,
              week_start, week_end. Used on the title slide.

    Returns:
        .pptx file bytes.
    """
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    blank_layout = prs.slide_layouts[6]
    meta = meta or {}

    for cs in slides:
        slide = prs.slides.add_slide(blank_layout)
        _fill_bg(slide, _LIGHT_BG)

        layout = cs.layout
        if layout == "title":
            _render_title(slide, cs, meta)
        elif layout == "concept_intro":
            _render_concept_intro(slide, cs)
        elif layout == "two_col":
            _render_two_col(slide, cs)
        elif layout == "visual_focus":
            _render_visual_focus(slide, cs)
        elif layout == "skill_card":
            _render_skill_card(slide, cs)
        elif layout == "summary":
            _render_summary(slide, cs)
        else:
            _render_two_col(slide, cs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Background ────────────────────────────────────────────────────────────────


def _fill_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Layout renderers ─────────────────────────────────────────────────────────


def _render_title(slide, cs: CompletedSlide, meta: dict) -> None:
    _add_rect(slide, 0, 0, _W, Inches(1.1), _BLUE)

    lesson_name = meta.get("lesson_name") or cs.title
    period_name = meta.get("period_name", "")
    grade = meta.get("grade_level", "")
    week_info = ""
    if meta.get("week_start"):
        week_info = f"Week of {meta['week_start']}"
        if meta.get("week_end"):
            week_info += f" – {meta['week_end']}"

    _add_text(
        slide, lesson_name,
        left=_MARGIN, top=Inches(0.15),
        width=_W - _MARGIN * 2, height=Inches(0.8),
        font_size=Pt(28), bold=True, colour=_WHITE, align=PP_ALIGN.LEFT,
    )

    sub_parts = [p for p in [period_name, grade, week_info] if p]
    sub_text = "  ·  ".join(sub_parts)
    _add_text(
        slide, sub_text,
        left=_MARGIN, top=Inches(1.25),
        width=_W - _MARGIN * 2, height=Inches(0.5),
        font_size=Pt(14), colour=_DARK, align=PP_ALIGN.LEFT,
    )

    if cs.bullets:
        _add_text(
            slide, "What we'll cover:",
            left=_MARGIN, top=Inches(2.0),
            width=_W - _MARGIN * 2, height=Inches(0.4),
            font_size=Pt(16), bold=True, colour=_BLUE, align=PP_ALIGN.LEFT,
        )
        bullet_text = "\n".join(f"• {b}" for b in cs.bullets)
        _add_text(
            slide, bullet_text,
            left=_MARGIN + Inches(0.2), top=Inches(2.5),
            width=_W - _MARGIN * 2, height=_H - Inches(2.8),
            font_size=Pt(15), colour=_DARK, align=PP_ALIGN.LEFT,
        )

    _add_rect(slide, 0, _H - Inches(0.12), _W, Inches(0.12), _ORANGE)
    _add_speaker_notes(slide, cs.speaker_notes)


def _render_concept_intro(slide, cs: CompletedSlide) -> None:
    _render_header_bar(slide, cs.title)

    bullets = list(cs.bullets)
    if bullets:
        # First bullet styled larger as the definition
        _add_text(
            slide, bullets[0],
            left=_MARGIN, top=_CONTENT_TOP,
            width=_W - _MARGIN * 2, height=Inches(1.4),
            font_size=Pt(20), bold=True, colour=_DARK, align=PP_ALIGN.LEFT,
        )
        if len(bullets) > 1:
            bullet_text = "\n".join(f"• {b}" for b in bullets[1:])
            _add_text(
                slide, bullet_text,
                left=_MARGIN, top=_CONTENT_TOP + Inches(1.5),
                width=_W - _MARGIN * 2, height=Inches(2.5),
                font_size=Pt(16), colour=_DARK, align=PP_ALIGN.LEFT,
            )

    if cs.prerequisites:
        prereq_y = _H - Inches(1.2)
        _add_rect(slide, _MARGIN, prereq_y, Inches(4.5), Inches(0.5), _ORANGE)
        prereq_label = "Prerequisites: " + ", ".join(cs.prerequisites)
        _add_text(
            slide, prereq_label,
            left=_MARGIN + Inches(0.1), top=prereq_y + Inches(0.05),
            width=Inches(4.3), height=Inches(0.4),
            font_size=Pt(11), bold=True, colour=_WHITE, align=PP_ALIGN.LEFT,
        )

    _add_speaker_notes(slide, cs.speaker_notes)


def _render_two_col(slide, cs: CompletedSlide) -> None:
    _render_header_bar(slide, cs.title)

    col_w = (_W - _MARGIN * 3) / 2
    text_left = _MARGIN
    img_left = _MARGIN * 2 + col_w
    content_h = _H - _CONTENT_TOP - _MARGIN

    bullet_text = "\n".join(f"• {b}" for b in cs.bullets)
    _add_text(
        slide, bullet_text,
        left=text_left, top=_CONTENT_TOP,
        width=col_w, height=content_h,
        font_size=Pt(15), colour=_DARK, align=PP_ALIGN.LEFT,
    )

    _place_visual(slide, cs, img_left, _CONTENT_TOP, col_w, content_h)
    _add_speaker_notes(slide, cs.speaker_notes)


def _render_visual_focus(slide, cs: CompletedSlide) -> None:
    _render_header_bar(slide, cs.title)

    img_h = _H - _CONTENT_TOP - Inches(1.5)
    caption_top = _CONTENT_TOP + img_h + Inches(0.1)

    _place_visual(slide, cs, _MARGIN, _CONTENT_TOP, _W - _MARGIN * 2, img_h)

    caption = cs.visual_caption or (cs.bullets[0] if cs.bullets else "")
    _add_text(
        slide, caption,
        left=_MARGIN, top=caption_top,
        width=_W - _MARGIN * 2, height=Inches(1.1),
        font_size=Pt(14), colour=_DARK, align=PP_ALIGN.CENTER,
    )

    _add_speaker_notes(slide, cs.speaker_notes)


def _render_skill_card(slide, cs: CompletedSlide) -> None:
    _render_header_bar(slide, cs.title)

    badge_y = _CONTENT_TOP
    next_x = _MARGIN

    if cs.bloom_level:
        badge_colour = _BLOOM_COLOURS.get(cs.bloom_level.lower(), _BLUE)
        _add_rect(slide, next_x, badge_y, Inches(2.2), Inches(0.45), badge_colour)
        _add_text(
            slide, f"Bloom: {cs.bloom_level}",
            left=next_x + Inches(0.05), top=badge_y + Inches(0.05),
            width=Inches(2.1), height=Inches(0.35),
            font_size=Pt(11), bold=True, colour=_WHITE, align=PP_ALIGN.CENTER,
        )
        next_x += Inches(2.4)

    if cs.difficulty:
        _add_rect(slide, next_x, badge_y, Inches(2.0), Inches(0.45), _ORANGE)
        _add_text(
            slide, f"Difficulty: {cs.difficulty}",
            left=next_x + Inches(0.05), top=badge_y + Inches(0.05),
            width=Inches(1.9), height=Inches(0.35),
            font_size=Pt(11), bold=True, colour=_WHITE, align=PP_ALIGN.CENTER,
        )

    badge_offset = Inches(0.7) if (cs.bloom_level or cs.difficulty) else Inches(0)

    _add_text(
        slide, "What mastery looks like:",
        left=_MARGIN, top=_CONTENT_TOP + badge_offset,
        width=_W - _MARGIN * 2, height=Inches(0.4),
        font_size=Pt(14), bold=True, colour=_BLUE, align=PP_ALIGN.LEFT,
    )

    bullet_text = "\n".join(f"• {b}" for b in cs.bullets)
    _add_text(
        slide, bullet_text,
        left=_MARGIN, top=_CONTENT_TOP + badge_offset + Inches(0.5),
        width=_W - _MARGIN * 2,
        height=_H - _CONTENT_TOP - badge_offset - Inches(0.5) - _MARGIN,
        font_size=Pt(15), colour=_DARK, align=PP_ALIGN.LEFT,
    )

    _add_speaker_notes(slide, cs.speaker_notes)


def _render_summary(slide, cs: CompletedSlide) -> None:
    _render_header_bar(slide, cs.title)

    bullet_text = "\n".join(f"✓  {b}" for b in cs.bullets)
    _add_text(
        slide, bullet_text,
        left=_MARGIN, top=_CONTENT_TOP,
        width=_W - _MARGIN * 2, height=_H - _CONTENT_TOP - _MARGIN,
        font_size=Pt(16), colour=_DARK, align=PP_ALIGN.LEFT,
    )

    _add_speaker_notes(slide, cs.speaker_notes)


# ── Shared helpers ────────────────────────────────────────────────────────────


def _render_header_bar(slide, title: str) -> None:
    _add_rect(slide, 0, 0, _W, Inches(1.1), _BLUE)
    _add_text(
        slide, title,
        left=_MARGIN, top=Inches(0.15),
        width=_W - _MARGIN * 2, height=Inches(0.8),
        font_size=Pt(24), bold=True, colour=_WHITE, align=PP_ALIGN.LEFT,
    )
    _add_rect(slide, 0, _H - Inches(0.08), _W, Inches(0.08), _ORANGE)


def _place_visual(
    slide,
    cs: CompletedSlide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
) -> None:
    """Place the approved/flagged image or a styled placeholder."""
    has_image = cs.visual_status in ("approved", "flagged") and cs.visual_path
    if has_image and os.path.exists(cs.visual_path):
        slide.shapes.add_picture(cs.visual_path, left, top, width=width, height=height)
    else:
        _add_rect(slide, left, top, width, height, _PLACEHOLDER_BG)
        label = (
            "Visual flagged for review"
            if cs.visual_status == "flagged"
            else "No visual"
        )
        _add_text(
            slide, label,
            left=left, top=top + height / 2 - Inches(0.3),
            width=width, height=Inches(0.6),
            font_size=Pt(13), colour=_MUTED, align=PP_ALIGN.CENTER,
        )


def _add_rect(slide, left, top, width, height, colour: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()


def _add_text(
    slide,
    text: str,
    left, top, width, height,
    font_size=Pt(14),
    bold: bool = False,
    colour: RGBColor = _DARK,
    align=PP_ALIGN.LEFT,
) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = colour


def _add_speaker_notes(slide, notes: str) -> None:
    if not notes:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes
