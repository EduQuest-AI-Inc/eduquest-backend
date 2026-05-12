import asyncio
import io
import pytest
from unittest.mock import MagicMock, AsyncMock, patch

from services.slides.pptx_generation_service import PptxGenerationService


def _make_pptx_bytes() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "[MOCK] Test Lesson"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _mock_agent(pptx_bytes=None, raises=None):
    """Mock PptxAgent: run() is AsyncMock returning {pptx_bytes, html_str} dict (or raising)."""
    agent = MagicMock()
    if raises:
        agent.run = AsyncMock(side_effect=raises)
    else:
        agent.run = AsyncMock(return_value={"pptx_bytes": pptx_bytes or _make_pptx_bytes(), "html_str": ""})
    return agent


def _svc(agent=None):
    provider = MagicMock()
    if agent is not None:
        provider.create_pptx_agent.return_value = agent
    svc = PptxGenerationService.__new__(PptxGenerationService)
    svc._bot_provider = provider
    svc.lesson_pptx_dao = MagicMock()
    svc.period_dao = MagicMock()
    return svc


def _curriculum(lesson_id="l1", lesson_name="Algebra Basics"):
    return {
        "lessons": [{"lesson_id": lesson_id, "lesson_name": lesson_name}],
        "concepts": [{"concept_name": "Variables", "lesson_name": lesson_name}],
        "skills": [{"skill_name": "Solve for x", "concept_name": "Variables"}],
        "concept_skills": [{"concept_name": "Variables", "skill_name": "Solve for x"}],
    }


def _period_context():
    return {
        "period_name": "Period 1",
        "grade_level": "9",
        "course_name": "Algebra I",
        "course_description": "Introduction to algebra",
    }


def _row(pptx_id="px1", lesson_id="l1", period_id="p1"):
    return {"pptx_id": pptx_id, "lesson_id": lesson_id, "period_id": period_id}


# ── _generate_one happy path ──────────────────────────────────────────────────

@pytest.mark.unit
def test_generate_one_happy_path():
    pptx_bytes = _make_pptx_bytes()
    svc = _svc(agent=_mock_agent(pptx_bytes=pptx_bytes))

    with patch("services.slides.pptx_generation_service.s3_service") as mock_s3:
        mock_s3.upload_pptx.return_value = "pptx/p1/l1.pptx"
        asyncio.run(
            svc._generate_one(_row(), _curriculum(), _period_context(), asyncio.Semaphore(1))
        )

    calls = svc.lesson_pptx_dao.update_status.call_args_list
    assert calls[0][0][1] == {"status": "generating"}
    assert calls[1][0][1]["status"] == "done"
    assert calls[1][0][1]["s3_key"] == "pptx/p1/l1.pptx"

    uploaded_bytes = mock_s3.upload_pptx.call_args[0][0]
    assert uploaded_bytes[:4] == b"PK\x03\x04", "expected valid PPTX (ZIP) magic bytes"


# ── _generate_one lesson not in curriculum ────────────────────────────────────

@pytest.mark.unit
def test_generate_one_lesson_not_found():
    svc = _svc(agent=_mock_agent())
    row = _row(lesson_id="missing")

    with patch("services.slides.pptx_generation_service.s3_service") as mock_s3:
        asyncio.run(
            svc._generate_one(row, _curriculum(lesson_id="l1"), _period_context(), asyncio.Semaphore(1))
        )

    svc.lesson_pptx_dao.update_status.assert_called_once_with("px1", {"status": "failed"})
    mock_s3.upload_pptx.assert_not_called()


# ── _generate_one agent raises ────────────────────────────────────────────────

@pytest.mark.unit
def test_generate_one_agent_raises():
    svc = _svc(agent=_mock_agent(raises=RuntimeError("agent exploded")))

    with patch("services.slides.pptx_generation_service.s3_service") as mock_s3:
        with pytest.raises(RuntimeError, match="agent exploded"):
            asyncio.run(
                svc._generate_one(_row(), _curriculum(), _period_context(), asyncio.Semaphore(1))
            )

    calls = svc.lesson_pptx_dao.update_status.call_args_list
    assert calls[0][0][1] == {"status": "generating"}
    assert len(calls) == 1  # exception surfaced before "failed" could be written
    mock_s3.upload_pptx.assert_not_called()


# ── _run_batch_async partial failure ─────────────────────────────────────────

@pytest.mark.unit
def test_run_batch_async_partial_failure():
    svc = _svc(agent=_mock_agent())
    rows = [
        _row(pptx_id="px1", lesson_id="l1"),
        _row(pptx_id="px2", lesson_id="missing"),  # not in curriculum → fails
    ]

    with patch("services.slides.pptx_generation_service.s3_service") as mock_s3:
        mock_s3.upload_pptx.return_value = "pptx/p1/l1.pptx"
        asyncio.run(svc._run_batch_async(rows, _curriculum(lesson_id="l1"), _period_context()))

    statuses = [c[0][1].get("status") for c in svc.lesson_pptx_dao.update_status.call_args_list]
    assert "done" in statuses, "l1 should reach done"
    assert "failed" in statuses, "missing lesson should reach failed"
