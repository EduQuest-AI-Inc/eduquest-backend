"""
Tests for the PPTX renderer. Uses mock CompletedSlide objects.
"""

from __future__ import annotations

import io

from pptx import Presentation

from models.slide_plan import CompletedSlide
from utils.rendering.pptx_renderer import render


def _slide(
    index: int,
    layout: str,
    title: str,
    bullets: list[str] | None = None,
    **extra,
) -> CompletedSlide:
    return CompletedSlide(
        index=index,
        layout=layout,
        title=title,
        bullets=bullets or ["Point one", "Point two"],
        speaker_notes="Teacher talking points here.",
        **extra,
    )


def test_render_returns_bytes():
    slides = [_slide(0, "title", "My Lesson")]
    result = render(slides)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_render_correct_slide_count():
    slides = [
        _slide(0, "title", "Intro"),
        _slide(1, "concept_intro", "Concept A"),
        _slide(2, "skill_card", "Skill A", bloom_level="Remember", difficulty="beginner"),
        _slide(3, "summary", "Wrap Up"),
    ]
    pptx_bytes = render(slides)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 4


def test_render_all_layouts():
    layouts = ["title", "concept_intro", "two_col", "visual_focus", "skill_card", "summary"]
    slides = [_slide(i, layout, f"Slide {i}") for i, layout in enumerate(layouts)]
    pptx_bytes = render(slides)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == len(layouts)


def test_render_two_col_with_placeholder():
    slide = _slide(0, "two_col", "Concept with Visual", visual_status="placeholder")
    pptx_bytes = render([slide])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 1


def test_render_flagged_visual_uses_placeholder():
    slide = _slide(0, "visual_focus", "Diagram Slide", visual_status="flagged")
    pptx_bytes = render([slide])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 1


def test_render_empty_slide_list():
    pptx_bytes = render([])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 0


def test_render_with_meta():
    slides = [_slide(0, "title", "Lesson Name")]
    meta = {
        "lesson_name": "Photosynthesis 101",
        "period_name": "Biology A",
        "grade_level": "9",
        "week_start": "2026-09-07",
        "week_end": "2026-09-11",
    }
    pptx_bytes = render(slides, meta=meta)
    assert isinstance(pptx_bytes, bytes)
    assert len(pptx_bytes) > 0


def test_render_skill_card_with_bloom_badges():
    slide = _slide(
        0,
        "skill_card",
        "Identify Cell Parts",
        bullets=["Label a cell diagram correctly", "Name 5 organelles"],
        bloom_level="Apply",
        difficulty="intermediate",
    )
    pptx_bytes = render([slide])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 1


def test_render_concept_intro_with_prereqs():
    slide = _slide(
        0,
        "concept_intro",
        "Cellular Respiration",
        bullets=["Conversion of glucose to ATP", "Occurs in mitochondria"],
        prerequisites=["ATP basics", "Cell anatomy"],
    )
    pptx_bytes = render([slide])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 1


def test_render_visual_focus_with_caption():
    slide = _slide(
        0,
        "visual_focus",
        "Cell Diagram",
        bullets=["A labeled view of a plant cell"],
        visual_caption="Plant cells contain chloroplasts that animal cells lack.",
        visual_status="placeholder",
    )
    pptx_bytes = render([slide])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 1


def test_render_with_real_image(tmp_path):
    """Render a two_col slide with an actual PNG image."""
    PNG_BYTES = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc"
        b"\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    img_path = tmp_path / "test.png"
    img_path.write_bytes(PNG_BYTES)

    slide = _slide(
        0,
        "two_col",
        "Concept with Real Image",
        visual_path=str(img_path),
        visual_status="approved",
    )
    pptx_bytes = render([slide])
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 1
