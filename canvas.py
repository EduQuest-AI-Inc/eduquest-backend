from canvasapi import Canvas
from bs4 import BeautifulSoup


class Course:
    def __init__(self, course_id, API_URL, API_KEY):
        self.API_URL = API_URL
        self.API_KEY = API_KEY
        self.canvas = Canvas(API_URL, API_KEY)
        self.course = self.canvas.get_course(course_id)
        self.course_id = course_id
        self.module_ids = [module.id for module in self.course.get_modules()]
        self.modules = []

    class Module:
        def __init__(self, Course, module_id):
            self.course = Course.course
            self.module = self.course.get_module(module_id)
            self.course_id = self.course.id
            self.title = self.module.name
            self.items = []
            self.item_ids = [item.id for item in self.module.get_module_items()]

        class Assignment:
            def __init__(self, assignment_id, course):
                self.course = course
                self.cv_assignment = self.course.get_assignment(assignment_id)
                self.title = self.cv_assignment.name
                if self.cv_assignment.description is not None:
                    html_content = self.cv_assignment.description
                    soup = BeautifulSoup(html_content, "html.parser")
                    for tag in soup(["script", "style", "link"]):
                        tag.decompose()
                    self.decription = soup.get_text(separator="\n", strip=True)
                self.due_date = self.cv_assignment.due_at_date

        class Quiz:  # do we want it?
            def __init__(self, quiz_id, course):
                self.course = course
                self.quiz = self.course.get_quiz(quiz_id)
                self.title = self.quiz.title
                if self.quiz.description is not None:
                    html_content = self.quiz.description
                    soup = BeautifulSoup(html_content, "html.parser")
                    for tag in soup(["script", "style", "link"]):
                        tag.decompose()
                    self.description = soup.get_text(separator="\n", strip=True)

        class File:
            def __init__(self, file_id, course):
                self.course = course
                self.file = self.course.get_file(file_id)
                self.id = file_id
                self.title = self.file.display_name
                self.mime_type = self.file.content_type
                self.download_url = self.file.url
                self.content = None
                self.processed_content = None
                self.metadata = {}

                # Download the file
                import requests
                response = requests.get(self.download_url)
                if response.status_code == 200:
                    self.content = response.content
                    self._process_file()

            def _process_file(self):
                """Process the file using AI-powered solutions to detect and convert content appropriately"""
                if not self.content:
                    return

                try:
                    # Use OpenAI's vision model for image analysis
                    if self._is_image():
                        self._process_image()
                    # Use OpenAI's document analysis for PDFs and documents
                    elif self._is_document():
                        self._process_document()
                    # Use OpenAI's code analysis for code files
                    elif self._is_code():
                        self._process_code()
                    # Use OpenAI's general text analysis for other text files
                    elif self._is_text():
                        self._process_text()
                    else:
                        print(f"Unsupported file type: {self.mime_type}")

                except Exception as e:
                    print(f"Error processing file {self.title}: {e}")

            def _is_image(self):
                """Check if the file is an image"""
                image_types = {
                    'image/jpeg', 'image/png', 'image/gif', 'image/bmp',
                    'image/webp', 'image/tiff', 'image/svg+xml'
                }
                return self.mime_type in image_types

            def _is_document(self):
                """Check if the file is a document"""
                document_types = {
                    'application/pdf',
                    'application/msword',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'application/vnd.ms-powerpoint',
                    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    'application/vnd.ms-excel'
                }
                return self.mime_type in document_types

            def _is_code(self):
                """Check if the file is a code file"""
                code_types = {
                    'text/x-python', 'text/x-java', 'text/x-c', 'text/x-c++',
                    'text/javascript', 'text/x-php', 'text/x-ruby', 'text/x-golang',
                    'text/x-csharp', 'application/typescript', 'text/x-tex'
                }
                return self.mime_type in code_types

            def _is_text(self):
                """Check if the file is a text file"""
                return self.mime_type.startswith('text/')

            def _process_image(self):
                """Process image files using OpenAI's vision model"""
                try:
                    import base64
                    from openai import OpenAI
                    client = OpenAI()

                    # Convert image to base64
                    base64_image = base64.b64encode(self.content).decode('utf-8')

                    # Use OpenAI's vision model to analyze the image
                    response = client.chat.completions.create(
                        model="gpt-4-vision-preview",
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text",
                                     "text": "Analyze this image and provide a detailed description of its content, including any text, diagrams, or important visual elements. Also identify any educational content or concepts being presented."},
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/jpeg;base64,{base64_image}"
                                        }
                                    }
                                ]
                            }
                        ],
                        max_tokens=1000
                    )

                    self.processed_content = response.choices[0].message.content
                    self.metadata['type'] = 'image_analysis'

                except Exception as e:
                    print(f"Error processing image: {e}")

            def _process_document(self):
                """Process document files using OpenAI's document analysis"""
                try:
                    from openai import OpenAI
                    client = OpenAI()

                    # Convert document content to text (preserving structure)
                    text_content = self._extract_text_from_document()

                    # Use OpenAI to analyze the document
                    response = client.chat.completions.create(
                        model="gpt-4",
                        messages=[
                            {
                                "role": "system",
                                "content": "You are an expert at analyzing educational documents. Extract key information, concepts, and structure from the document while preserving important formatting and relationships between elements."
                            },
                            {
                                "role": "user",
                                "content": text_content
                            }
                        ],
                        max_tokens=2000
                    )

                    self.processed_content = response.choices[0].message.content
                    self.metadata['type'] = 'document_analysis'

                except Exception as e:
                    print(f"Error processing document: {e}")

            def _process_code(self):
                """Process code files using OpenAI's code analysis"""
                try:
                    from openai import OpenAI
                    client = OpenAI()

                    # Decode code content
                    code_content = self.content.decode('utf-8')

                    # Use OpenAI to analyze the code
                    response = client.chat.completions.create(
                        model="gpt-4",
                        messages=[
                            {
                                "role": "system",
                                "content": "You are an expert at analyzing educational code. Explain the code's purpose, structure, and key concepts while preserving important implementation details."
                            },
                            {
                                "role": "user",
                                "content": code_content
                            }
                        ],
                        max_tokens=2000
                    )

                    self.processed_content = response.choices[0].message.content
                    self.metadata['type'] = 'code_analysis'

                except Exception as e:
                    print(f"Error processing code: {e}")

            def _process_text(self):
                """Process text files using OpenAI's text analysis"""
                try:
                    from openai import OpenAI
                    client = OpenAI()

                    # Decode text content
                    text_content = self.content.decode('utf-8')

                    # Use OpenAI to analyze the text
                    response = client.chat.completions.create(
                        model="gpt-4",
                        messages=[
                            {
                                "role": "system",
                                "content": "You are an expert at analyzing educational text. Extract key information, concepts, and structure while preserving important relationships between elements."
                            },
                            {
                                "role": "user",
                                "content": text_content
                            }
                        ],
                        max_tokens=2000
                    )

                    self.processed_content = response.choices[0].message.content
                    self.metadata['type'] = 'text_analysis'

                except Exception as e:
                    print(f"Error processing text: {e}")

            def _extract_text_from_document(self):
                """Extract text from various document types while preserving structure"""
                try:
                    if self.mime_type == 'application/pdf':
                        import PyPDF2
                        from io import BytesIO
                        pdf_file = BytesIO(self.content)
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        text = ""
                        for page in pdf_reader.pages:
                            text += page.extract_text() + "\n"
                        return text

                    elif self.mime_type in ['application/msword',
                                            'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                        import docx
                        from io import BytesIO
                        doc_file = BytesIO(self.content)
                        doc = docx.Document(doc_file)
                        text = ""
                        for paragraph in doc.paragraphs:
                            text += paragraph.text + "\n"
                        return text

                    elif self.mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                        import pptx
                        from io import BytesIO
                        ppt_file = BytesIO(self.content)
                        ppt = pptx.Presentation(ppt_file)
                        text = ""
                        for slide in ppt.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                        return text

                    return None

                except Exception as e:
                    print(f"Error extracting text from document: {e}")
                    return None

        class Page:
            def __init__(self, page_id, course):
                self.course = course
                self.page = self.course.get_page(page_id)
                self.title = self.page.title
                self.body = None
                if self.page.body is not None:
                    html_content = self.page.body
                    soup = BeautifulSoup(html_content, "html.parser")
                    for tag in soup(["script", "style", "link"]):
                        tag.decompose()
                    self.body = soup.get_text(separator="\n", strip=True)

        class Discussion:
            def __init__(self, discussion_id, course):
                self.course = course
                self.discussion = self.course.get_discussion_topic(discussion_id)
                self.title = self.discussion.title
                if self.discussion.message is not None:
                    html_content = self.discussion.message
                    soup = BeautifulSoup(html_content, "html.parser")
                    for tag in soup(["script", "style", "link"]):
                        tag.decompose()
                    self.body = soup.get_text(separator="\n", strip=True)

        def get_item(self):
            for id in self.item_ids:
                item = self.module.get_module_item(id)
                if item.type == 'Assignment':
                    item = self.Assignment(item.content_id, self.course)
                elif item.type == 'Quiz':
                    item = self.Quiz(item.content_id, self.course)
                elif item.type == 'File':
                    item = self.File(item.content_id, self.course)
                elif item.type == 'Page':
                    item = self.Page(item.content_id, self.course)
                elif item.type == 'Discussion':
                    item = self.Discussion(item.content_id, self.course)
                else:
                    pass

            self.items.append(item)
            # elif item.type == 'Text Header':
            #
            # elif item.type == 'External URL':
            #     pass
            # elif item.type == 'External Tool':
            #     pass

    # def retrieve_pages(self):
    #     pages = self.course.get_pages()
    #     ids = [page.page_id for page in self.course.get_pages()]
    #     page_data = []
    #     for id in ids:
    #         page_info = {
    #             "id": self.course.get_page(id).page_id,
    #             "title": self.course.get_page(id).title,
    #             "body": None
    #         }
    #         if self.course.get_page(id).body is not None:
    #             html_content = self.course.get_page(id).body
    #             soup = BeautifulSoup(html_content, "html.parser")
    #             for tag in soup(["script", "style", "link"]):
    #                 tag.decompose()
    #             page_info["body"] = soup.get_text(separator="\n", strip=True)
    #         page_data.append(page_info)
    #     return page_data

    # def retrieve_files(self):
    #     pass
    #
    # def retrieve_modules(self):
    #     pass
    #
    # def retrieve_assignments(self):
    #     assignments = self.course.get_assignments()
    #     assignment_data = []
    #     for assignment in assignments:
    #         assignment_info = {
    #             "id": assignment.id,
    #             "name": assignment.name,
    #             "description": None
    #         }
    #         if assignment.description is not None:
    #             html_content = assignment.description
    #             soup = BeautifulSoup(html_content, "html.parser")
    #             for tag in soup(["script", "style", "link"]):
    #                 tag.decompose()
    #             assignment_info["description"] =
    #         assignment_data.append(assignment_info)
    #     return assignment_data
    #
    # def retrieve_quizzes(self):
    #     pass
    #
    # def retrieve_syllabus(self):
    #     pass

