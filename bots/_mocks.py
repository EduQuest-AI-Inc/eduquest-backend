"""
Mock bot implementations for fast testing without OpenAI API calls.

Activated when MOCK_AI=true via bots/provider.py. All classes return
realistic-shaped data (real Pydantic types) so downstream service logic
runs unchanged.
"""
import math
from datetime import date
from typing import Any


class MockRunResult:
    """Duck-typed stand-in for the result returned by agents.Runner.run()."""

    def __init__(self, final_output: Any, last_response_id: str = "mock-response-id-001"):
        self.final_output = final_output
        self.last_response_id = last_response_id


class MockRunner:
    """
    Drop-in replacement for agents.Runner. Dispatches on agent.output_type
    so the right Pydantic response model is returned for each conversation agent.

    ProfileResponse turns (tracked via previous_response_id):
      None / unknown → turn 1: ask about subjects          → mock-response-id-001
      mock-response-id-001 → turn 2: follow-up question    → mock-response-id-002
      mock-response-id-002 → turn 3: complete with profile → mock-response-id-003
    """

    @staticmethod
    async def run(agent: Any, message: str, **kwargs) -> MockRunResult:
        from bots.ltg_agent import LTGResponse
        from bots.profile_agent import ProfileResponse, StudentProfile
        from bots.teacher_feedback_agent import TeacherFeedbackResponse

        output_type = getattr(agent, "output_type", None)

        if output_type is LTGResponse:
            prev = kwargs.get("previous_response_id")
            if prev == "mock-ltg-response-id-001":
                # Turn 2: student picked a goal — confirm and close the conversation
                ltg_payload = LTGResponse(
                    message="[MOCK] Excellent choice! I've selected 'Build a portfolio project that applies core course concepts' as your long-term goal.",
                    goal_1=None,
                    goal_2=None,
                    goal_3=None,
                    chosen_goal="Build a portfolio project that applies core course concepts",
                )
                return MockRunResult(final_output=ltg_payload, last_response_id="mock-ltg-response-id-002")
            # Turn 1 (initiate) — no previous_response_id
            ltg_payload = LTGResponse(
                message="[MOCK] Here are three long-term goal options based on your course. Which one resonates with you?",
                goal_1="Build a portfolio project that applies core course concepts",
                goal_2="Teach a concept you learned to a peer or family member",
                goal_3="Apply a skill from class to solve a real-world problem you care about",
                chosen_goal=None,
            )
            return MockRunResult(final_output=ltg_payload, last_response_id="mock-ltg-response-id-001")

        if output_type is ProfileResponse:
            prev = kwargs.get("previous_response_id")
            if prev == "mock-response-id-001":
                profile_payload = ProfileResponse(
                    response="[MOCK] That's great! Tell me about something you find challenging in school and how you usually like to learn new things.",
                    profile=None,
                )
                return MockRunResult(final_output=profile_payload, last_response_id="mock-response-id-002")
            if prev == "mock-response-id-002":
                profile_payload = ProfileResponse(
                    response="[MOCK] Thanks for sharing! I have everything I need to personalize your experience. Let's get started!",
                    profile=StudentProfile(
                        strengths=["Analytical thinking", "Problem-solving"],
                        weaknesses=["Time management", "Public speaking"],
                        interests=["Technology", "Science"],
                        learning_styles=["Visual", "Hands-on"],
                    ),
                )
                return MockRunResult(final_output=profile_payload, last_response_id="mock-response-id-003")
            # Turn 1 (initiate) — no previous_response_id
            profile_payload = ProfileResponse(
                response="[MOCK] Hi! I'm EduQuest. I'd love to learn more about you. What subjects do you enjoy most?",
                profile=None,
            )
            return MockRunResult(final_output=profile_payload, last_response_id="mock-response-id-001")

        if output_type is TeacherFeedbackResponse:
            feedback_payload = TeacherFeedbackResponse(
                response="[MOCK] Based on the quest data, this student shows strong engagement with hands-on tasks and may benefit from more open-ended challenges.",
                suggested_change=None,
            )
            return MockRunResult(final_output=feedback_payload)

        raise ValueError(f"MockRunner: unrecognized agent output_type {output_type!r}")


class MockHWAgent:
    """Fast replacement for HWAgent — returns one IndividualQuest per schedule item."""

    def __init__(self, student, period, schedule, conversation_id=None, previous_response_id=None):
        self.schedule = schedule

    def run(self) -> list:
        from bots.quests.quest_agent import IndividualQuest

        results = []
        for quest in self.schedule:
            name = quest.get("Name", "Quest") if isinstance(quest, dict) else getattr(quest, "Name", "Quest")
            skills = quest.get("Skills", "") if isinstance(quest, dict) else getattr(quest, "Skills", "")
            week = quest.get("Week", 1) if isinstance(quest, dict) else getattr(quest, "Week", 1)
            results.append(IndividualQuest(
                Name=f"[MOCK] {name}",
                Skills=skills,
                Week=week,
                instructions=[
                    {"step": 1, "text": "[MOCK] Review the week's material carefully."},
                    {"step": 2, "text": "Complete the main activity described in class."},
                    {"step": 3, "text": "Reflect on what you learned and submit your work."},
                ],
                rubric={
                    "Criteria": {
                        "Understanding": {
                            "Score_0": "No attempt",
                            "Score_1": "Minimal understanding",
                            "Score_2": "Partial understanding",
                            "Score_3": "Satisfactory understanding",
                            "Score_4": "Good understanding",
                            "Score_5": "Excellent understanding",
                        }
                    }
                },
            ))
        return results


class MockLTGScheduleAgent:
    """Fast replacement for LTGScheduleAgent — returns one goal-aligned name per schedule item."""

    def __init__(self, student=None, schedule=None, goal_text=None, **kwargs):
        self._schedule = schedule or []
        self._goal_text = goal_text or "complete the course"

    def run(self):
        from bots.quests.ltg_schedule_agent import ScheduleOutput, WeekQuest

        verbs = ["Build", "Design", "Analyze", "Create", "Apply",
                 "Investigate", "Prototype", "Compare", "Draft", "Evaluate"]
        goal_fragment = self._goal_text[:40].rstrip()
        quests = []
        for i, entry in enumerate(self._schedule):
            week = entry.get("Week", i + 1) if isinstance(entry, dict) else i + 1
            skills_raw = entry.get("Skills", "") if isinstance(entry, dict) else ""
            first_skill = skills_raw.split(";")[0].strip() if skills_raw else "core skills"
            verb = verbs[i % len(verbs)]
            quest_name = f"[MOCK] {verb} a {first_skill} artifact toward: {goal_fragment}"
            quests.append(WeekQuest(week=week, quest_name=quest_name))
        return ScheduleOutput(quests=quests)


class MockGradingOrchestrator:
    """Fast replacement for GradingOrchestrator."""

    async def grade_submission(self, grading_input: Any) -> Any:
        from bots.grading_agent import GradingResult

        skill_mastery = {skill: 0.75 for skill in (grading_input.skills or ["general"])}
        return GradingResult(
            numerical_grade=38,
            feedback=(
                "[MOCK] Good effort! Your submission demonstrates solid understanding of the core concepts. "
                "Consider expanding your analysis with more specific examples in future work."
            ),
            skill_mastery=skill_mastery,
            homework_changes_recommended=False,
            recommended_changes=None,
        )


_RESEARCH_TOPICS = [
    "Research-based Topic A",
    "Research-based Topic B",
    "Research-based Topic C",
]

_MIN_KEYWORD_LENGTH = 4
_MAX_KEYWORDS = 5

_BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"]
_DIFFICULTIES = ["beginner", "intermediate", "advanced"]


def _extract_topics(course_description: str) -> list[str]:
    """Derive up to 5 keywords from course_description, falling back to research placeholders."""
    words = [w for w in course_description.split() if len(w) >= _MIN_KEYWORD_LENGTH]
    return words[:_MAX_KEYWORDS] if words else _RESEARCH_TOPICS


class MockCurriculumAgent:
    """
    Deterministic replacement for CurriculumAgent. Returns a fully structured
    CurriculumResult with no randomness, no network calls, and no DB writes.
    The service layer is responsible for persisting the result.

    Determinism rules:
      - Week count  = ceil((end_date - start_date).days / 7)
      - Lessons     = 1 if week_number is odd, 2 if even
      - Concepts    = 1 if lesson_index is odd (1-based), 2 if even
      - Skills      = 1 if concept_index is odd (1-based), 2 if even
      - Topics      = keywords cycled from course_description (or research placeholders)
    """

    def __init__(
        self,
        vector_store_ids: list,
        course_name: str | None = None,
        start_date: str | None = None,
        end_date: str | None = None,
        course_description: str | None = None,
        grade_level: str | None = None,
        research_context: str | None = None,
    ):
        self._start_date = date.fromisoformat(start_date) if start_date else date.today()
        self._end_date = date.fromisoformat(end_date) if end_date else date.today()
        self._course_description = course_description or course_name or ""

    def run(self):
        from bots.schemas.curriculum import (
            CurriculumConcept,
            CurriculumLesson,
            CurriculumResult,
            CurriculumSkill,
            CurriculumWeek,
        )

        total_weeks = math.ceil((self._end_date - self._start_date).days / 7) or 3
        topics = _extract_topics(self._course_description)
        weeks = []

        for w in range(1, total_weeks + 1):
            topic = topics[(w - 1) % len(topics)]
            lesson_count = 1 if w % 2 != 0 else 2
            lessons = []

            for lesson_num in range(1, lesson_count + 1):
                concept_count = 1 if lesson_num % 2 != 0 else 2
                concepts = []

                for c in range(1, concept_count + 1):
                    skill_count = 1 if c % 2 != 0 else 2
                    skills = []
                    for s in range(1, skill_count + 1):
                        idx = w + lesson_num + c + s
                        skills.append(CurriculumSkill(
                            skill_id=f"{w}.{lesson_num}.{c}.{s}",
                            title=f"Skill {w}.{lesson_num}.{c}.{s}",
                            description=f"Students can demonstrate {_BLOOM_LEVELS[(idx) % len(_BLOOM_LEVELS)].lower()}-level mastery of skill {w}.{lesson_num}.{c}.{s}.",
                            bloom_level=_BLOOM_LEVELS[idx % len(_BLOOM_LEVELS)],
                            difficulty=_DIFFICULTIES[(w + s) % len(_DIFFICULTIES)],
                            mastery_threshold=0.8 if s % 2 != 0 else 0.7,
                        ))

                    cidx = w + lesson_num + c
                    prereqs = [f"Familiarity with Week {w - 1} material"] if w > 1 else []
                    if c > 1:
                        prereqs.append(f"Concept {w}.{lesson_num}.{c - 1}")
                    concepts.append(
                        CurriculumConcept(
                            concept_id=f"{w}.{lesson_num}.{c}",
                            title=f"Concept {w}.{lesson_num}.{c}",
                            description=f"An exploration of concept {w}.{lesson_num}.{c} within {topic}, focusing on core principles and real-world application.",
                            prerequisites=prereqs,
                            key_takeaways=[
                                f"Understand the foundational principles of concept {w}.{lesson_num}.{c}.",
                                f"Apply concept {w}.{lesson_num}.{c} to solve problems in {topic}.",
                            ] + ([f"Evaluate trade-offs when using concept {w}.{lesson_num}.{c}."] if cidx % 2 == 0 else []),
                            common_misconceptions=(
                                [f"Students often confuse concept {w}.{lesson_num}.{c} with adjacent ideas from earlier weeks."]
                                if cidx % 3 != 0 else []
                            ),
                            skills=skills,
                        )
                    )

                lessons.append(
                    CurriculumLesson(
                        lesson_id=f"{w}.{lesson_num}",
                        title=f"Lesson {w}.{lesson_num}",
                        concepts=concepts,
                    )
                )

            weeks.append(
                CurriculumWeek(
                    week_number=w,
                    week_id=str(w),
                    title=f"Week {w}: {topic}",
                    lessons=lessons,
                )
            )

        return CurriculumResult(
            grade_level="unspecified",
            course=self._course_description,
            total_weeks=total_weeks,
            weeks=weeks,
        )

    async def run_async(self):
        return self.run()

    def run_and_get_json(self) -> dict:
        return self.run().model_dump()

    async def run_and_get_json_async(self) -> dict:
        return self.run().model_dump()


class MockConversationsSession:
    """Duck-typed stand-in for agents.OpenAIConversationsSession."""

    def __init__(self, conversation_id=None):
        self._session_id = conversation_id or "mock-conversation-id-001"


class MockPptxAgent:
    """
    Deterministic replacement for PptxAgent. Returns a minimal but real
    .pptx file (title slide + one content slide per concept) so the file
    can be opened after downloading. Uses python-pptx.
    """

    async def run(self, lesson: dict, period_context: dict) -> dict:
        import io
        from pptx import Presentation

        prs = Presentation()
        lesson_name = lesson.get("lesson_name", "Lesson")

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = f"[MOCK] {lesson_name}"
        title_slide.placeholders[1].text = "EduQuest — Mock PowerPoint"

        concepts = lesson.get("concepts", [])
        skills = lesson.get("skills", [])
        for concept in concepts:
            concept_name = concept.get("concept_name", "Concept")
            concept_skills = [s for s in skills if s.get("concept_name") == concept_name]
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = concept_name
            tf = slide.placeholders[1].text_frame
            tf.text = "Skills:"
            for skill in concept_skills:
                p = tf.add_paragraph()
                p.text = f"• {skill.get('skill_name', '')}"
                p.level = 1

        buf = io.BytesIO()
        prs.save(buf)
        return {"pptx_bytes": buf.getvalue(), "html_str": "<html><body>[MOCK]</body></html>"}
